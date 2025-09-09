import streamlit as st
import json
import os
import urllib.request
import urllib.parse
from datetime import datetime
from typing import Dict, List, Any
import re
import time
import hashlib
import csv

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    from fake_useragent import UserAgent
    UA_AVAILABLE = True
except ImportError:
    UA_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class SimpleFakeUserAgent:
    def __init__(self):
        self.agents = [
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:89.0) Gecko/20100101 Firefox/89.0',
            'Mozilla/5.0 (Macintosh; Intel Mac OS X 10.15; rv:89.0) Gecko/20100101 Firefox/89.0'
        ]
        self.current_index = 0
    
    @property
    def random(self):
        agent = self.agents[self.current_index]
        self.current_index = (self.current_index + 1) % len(self.agents)
        return agent

class DocumentExtractor:
    def __init__(self):
        self.supported_formats = [
            '.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.txt', '.csv'
        ]
    
    def extract_text(self, file_path: str, file_content: bytes = None) -> List[Dict[str, Any]]:
        try:
            file_extension = os.path.splitext(file_path.lower())[1]
            
            if file_extension == '.pdf':
                return self._extract_pdf(file_path, file_content)
            elif file_extension == '.docx':
                return self._extract_docx(file_path, file_content)
            elif file_extension in ['.xlsx', '.xls']:
                return self._extract_excel(file_path, file_content)
            elif file_extension == '.pptx':
                return self._extract_pptx(file_path, file_content)
            elif file_extension == '.txt':
                return self._extract_txt(file_path, file_content)
            elif file_extension == '.csv':
                return self._extract_csv(file_path, file_content)
            else:
                return self._extract_generic_text(file_path, file_content)
                
        except Exception as e:
            st.error(f"Error extracting from {file_path}: {e}")
            return []
    
    def _extract_pdf(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        extracted_data = []
        
        if PDFPLUMBER_AVAILABLE and file_content:
            try:
                import io
                with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text and text.strip():
                            extracted_data.append({
                                'id': hashlib.md5(f"{file_path}_{page_num}_{text[:100]}".encode()).hexdigest()[:12],
                                'type': 'document_page',
                                'platform': 'PDF Document',
                                'title': f'Page {page_num + 1}',
                                'content': text.strip(),
                                'author': 'Document',
                                'url': file_path,
                                'timestamp': datetime.now().isoformat(),
                                'engagement': len(text.split()),
                                'post_type': 'pdf_page'
                            })
                return extracted_data
            except Exception as e:
                st.warning(f"pdfplumber failed: {e}, trying PyPDF2")
        
        if PYPDF2_AVAILABLE and file_content:
            try:
                import io
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        extracted_data.append({
                            'id': hashlib.md5(f"{file_path}_{page_num}_{text[:100]}".encode()).hexdigest()[:12],
                            'type': 'document_page',
                            'platform': 'PDF Document',
                            'title': f'Page {page_num + 1}',
                            'content': text.strip(),
                            'author': 'Document',
                            'url': file_path,
                            'timestamp': datetime.now().isoformat(),
                            'engagement': len(text.split()),
                            'post_type': 'pdf_page'
                        })
                return extracted_data
            except Exception as e:
                st.warning(f"PyPDF2 failed: {e}")
        
        return self._extract_generic_text(file_path, file_content)
    
    def _extract_docx(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        if not DOCX_AVAILABLE:
            return self._extract_generic_text(file_path, file_content)
        
        try:
            import io
            doc = Document(io.BytesIO(file_content))
            extracted_data = []
            
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text.strip())
            
            content = '\n'.join(full_text)
            if content:
                extracted_data.append({
                    'id': hashlib.md5(f"{file_path}_{content[:100]}".encode()).hexdigest()[:12],
                    'type': 'document',
                    'platform': 'DOCX Document',
                    'title': os.path.basename(file_path),
                    'content': content,
                    'author': 'Document',
                    'url': file_path,
                    'timestamp': datetime.now().isoformat(),
                    'engagement': len(content.split()),
                    'post_type': 'docx_document'
                })
            
            return extracted_data
            
        except Exception as e:
            st.warning(f"DOCX extraction failed: {e}")
            return self._extract_generic_text(file_path, file_content)
    
    def _extract_excel(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        extracted_data = []
        
        if PANDAS_AVAILABLE:
            try:
                import io
                if file_path.lower().endswith('.xlsx'):
                    df = pd.read_excel(io.BytesIO(file_content), sheet_name=None)
                else:
                    df = pd.read_excel(io.BytesIO(file_content), sheet_name=None, engine='xlrd')
                
                for sheet_name, sheet_df in df.items():
                    text_content = []
                    for column in sheet_df.columns:
                        text_content.append(f"{column}: {', '.join(sheet_df[column].astype(str).tolist())}")
                    
                    content = '\n'.join(text_content)
                    if content:
                        extracted_data.append({
                            'id': hashlib.md5(f"{file_path}_{sheet_name}_{content[:100]}".encode()).hexdigest()[:12],
                            'type': 'spreadsheet',
                            'platform': 'Excel Document',
                            'title': f'Sheet: {sheet_name}',
                            'content': content,
                            'author': 'Document',
                            'url': file_path,
                            'timestamp': datetime.now().isoformat(),
                            'engagement': len(content.split()),
                            'post_type': 'excel_sheet'
                        })
                
                return extracted_data
                
            except Exception as e:
                st.warning(f"Excel extraction with pandas failed: {e}")
        
        if OPENPYXL_AVAILABLE and file_path.lower().endswith('.xlsx'):
            try:
                import io
                wb = openpyxl.load_workbook(io.BytesIO(file_content))
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    text_content = []
                    for row in sheet.iter_rows(values_only=True):
                        row_text = ', '.join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text_content.append(row_text)
                    
                    content = '\n'.join(text_content)
                    if content:
                        extracted_data.append({
                            'id': hashlib.md5(f"{file_path}_{sheet_name}_{content[:100]}".encode()).hexdigest()[:12],
                            'type': 'spreadsheet',
                            'platform': 'Excel Document',
                            'title': f'Sheet: {sheet_name}',
                            'content': content,
                            'author': 'Document',
                            'url': file_path,
                            'timestamp': datetime.now().isoformat(),
                            'engagement': len(content.split()),
                            'post_type': 'excel_sheet'
                        })
                
                return extracted_data
                
            except Exception as e:
                st.warning(f"Excel extraction with openpyxl failed: {e}")
        
        return self._extract_generic_text(file_path, file_content)
    
    def _extract_pptx(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        if not PPTX_AVAILABLE:
            return self._extract_generic_text(file_path, file_content)
        
        try:
            import io
            prs = Presentation(io.BytesIO(file_content))
            extracted_data = []
            
            for slide_num, slide in enumerate(prs.slides):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                content = '\n'.join(text_content)
                if content:
                    extracted_data.append({
                        'id': hashlib.md5(f"{file_path}_{slide_num}_{content[:100]}".encode()).hexdigest()[:12],
                        'type': 'presentation_slide',
                        'platform': 'PowerPoint Document',
                        'title': f'Slide {slide_num + 1}',
                        'content': content,
                        'author': 'Document',
                        'url': file_path,
                        'timestamp': datetime.now().isoformat(),
                        'engagement': len(content.split()),
                        'post_type': 'pptx_slide'
                    })
            
            return extracted_data
            
        except Exception as e:
            st.warning(f"PowerPoint extraction failed: {e}")
            return self._extract_generic_text(file_path, file_content)
    
    def _extract_txt(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        try:
            content = file_content.decode('utf-8')
            if content.strip():
                return [{
                    'id': hashlib.md5(f"{file_path}_{content[:100]}".encode()).hexdigest()[:12],
                    'type': 'text_file',
                    'platform': 'Text Document',
                    'title': os.path.basename(file_path),
                    'content': content.strip(),
                    'author': 'Document',
                    'url': file_path,
                    'timestamp': datetime.now().isoformat(),
                    'engagement': len(content.split()),
                    'post_type': 'text_file'
                }]
        except UnicodeDecodeError:
            try:
                content = file_content.decode('latin-1')
                if content.strip():
                    return [{
                        'id': hashlib.md5(f"{file_path}_{content[:100]}".encode()).hexdigest()[:12],
                        'type': 'text_file',
                        'platform': 'Text Document',
                        'title': os.path.basename(file_path),
                        'content': content.strip(),
                        'author': 'Document',
                        'url': file_path,
                        'timestamp': datetime.now().isoformat(),
                        'engagement': len(content.split()),
                        'post_type': 'text_file'
                    }]
            except Exception as e:
                st.error(f"Text file encoding error: {e}")
        return []
    
    def _extract_csv(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        try:
            content = file_content.decode('utf-8')
            lines = content.split('\n')
            text_content = []
            
            for line_num, line in enumerate(lines):
                if line.strip():
                    text_content.append(f"Row {line_num + 1}: {line.strip()}")
            
            full_content = '\n'.join(text_content)
            if full_content:
                return [{
                    'id': hashlib.md5(f"{file_path}_{full_content[:100]}".encode()).hexdigest()[:12],
                    'type': 'csv_file',
                    'platform': 'CSV Document',
                    'title': os.path.basename(file_path),
                    'content': full_content,
                    'author': 'Document',
                    'url': file_path,
                    'timestamp': datetime.now().isoformat(),
                    'engagement': len(full_content.split()),
                    'post_type': 'csv_file'
                }]
        except Exception as e:
            st.error(f"CSV extraction error: {e}")
        return []
    
    def _extract_generic_text(self, file_path: str, file_content: bytes) -> List[Dict[str, Any]]:
        try:
            content = file_content.decode('utf-8', errors='ignore')
            if content.strip():
                return [{
                    'id': hashlib.md5(f"{file_path}_{content[:100]}".encode()).hexdigest()[:12],
                    'type': 'generic_file',
                    'platform': 'Generic Document',
                    'title': os.path.basename(file_path),
                    'content': content.strip(),
                    'author': 'Document',
                    'url': file_path,
                    'timestamp': datetime.now().isoformat(),
                    'engagement': len(content.split()),
                    'post_type': 'generic_file'
                }]
        except Exception as e:
            st.error(f"Generic text extraction error: {e}")
        return []

class DataProcessor:
    def __init__(self):
        self.standard_fields = [
            'id', 'type', 'platform', 'title', 'content', 'author', 
            'url', 'timestamp', 'scraped_at', 'engagement', 'post_type'
        ]
    
    def process(self, raw_data: List[Dict[str, Any]], include_metadata: bool = True) -> List[Dict[str, Any]]:
        processed_data = []
        
        for item in raw_data:
            processed_item = self._standardize_item(item, include_metadata)
            if processed_item:
                processed_data.append(processed_item)
        
        return processed_data
    
    def _standardize_item(self, item: Dict[str, Any], include_metadata: bool) -> Dict[str, Any]:
        try:
            content_hash = hashlib.md5(
                f"{item.get('content', '')}{item.get('url', '')}{item.get('timestamp', '')}".encode()
            ).hexdigest()[:12]
            
            standardized = {
                'id': content_hash,
                'type': item.get('type', 'unknown'),
                'platform': item.get('platform', 'unknown'),
                'title': self._clean_text(item.get('title', '')),
                'content': self._clean_text(item.get('content', '')),
                'author': self._clean_text(item.get('author', 'Unknown')),
                'url': item.get('url', ''),
                'timestamp': self._standardize_timestamp(item.get('timestamp') or item.get('created_utc')),
                'scraped_at': datetime.now().isoformat(),
                'engagement': self._extract_numeric_value(item.get('engagement') or item.get('score', 0)),
                'post_type': item.get('post_type', 'general')
            }
            
            if include_metadata:
                standardized.update({
                    'word_count': len(standardized['content'].split()) if standardized['content'] else 0,
                    'char_count': len(standardized['content']) if standardized['content'] else 0,
                    'has_title': bool(standardized['title']),
                    'content_length_category': self._categorize_content_length(standardized['content'])
                })
                
                for key, value in item.items():
                    if key not in self.standard_fields and not key.startswith('_'):
                        standardized[f'extra_{key}'] = value
            
            return standardized
            
        except Exception as e:
            print(f"Error processing item: {e}")
            return None
    
    def _clean_text(self, text: str) -> str:
        if not text:
            return ''
        
        cleaned = ' '.join(text.split())
        cleaned = cleaned.replace('\r', ' ').replace('\n', ' ')
        
        return cleaned.strip()
    
    def _standardize_timestamp(self, timestamp) -> str:
        if not timestamp:
            return ''
        
        try:
            if isinstance(timestamp, (int, float)):
                return datetime.fromtimestamp(timestamp).isoformat()
            
            if isinstance(timestamp, str):
                formats = [
                    '%Y-%m-%d %H:%M:%S',
                    '%Y-%m-%dT%H:%M:%S',
                    '%Y-%m-%dT%H:%M:%SZ',
                    '%Y-%m-%d',
                    '%m/%d/%Y',
                    '%d/%m/%Y'
                ]
                
                for fmt in formats:
                    try:
                        return datetime.strptime(timestamp, fmt).isoformat()
                    except ValueError:
                        continue
                
                return timestamp
            
            return str(timestamp)
            
        except Exception:
            return str(timestamp) if timestamp else ''
    
    def _extract_numeric_value(self, value) -> int:
        if isinstance(value, (int, float)):
            return int(value)
        
        if isinstance(value, str):
            numbers = re.findall(r'\d+', value.replace(',', ''))
            return int(numbers[0]) if numbers else 0
        
        return 0
    
    def _categorize_content_length(self, content: str) -> str:
        if not content:
            return 'empty'
        
        length = len(content)
        if length < 50:
            return 'short'
        elif length < 200:
            return 'medium'
        elif length < 500:
            return 'long'
        else:
            return 'very_long'

class FileManager:
    def __init__(self):
        self.data_dir = "data"
        self._ensure_data_directory()
    
    def _ensure_data_directory(self):
        if not os.path.exists(self.data_dir):
            os.makedirs(self.data_dir)
    
    def save_data(self, data: List[Dict[str, Any]], format_type: str, filename: str) -> str:
        if not data:
            raise ValueError("No data to save")
        
        if not filename:
            filename = "social_media_data"
        
        filename = filename.split('.')[0]
        
        if format_type == 'json':
            return self._save_json(data, filename)
        elif format_type == 'csv':
            return self._save_csv(data, filename)
        elif format_type == 'excel':
            return self._save_excel(data, filename)
        else:
            raise ValueError(f"Unsupported format: {format_type}")
    
    def _save_json(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        filepath = os.path.join(self.data_dir, f"{filename}.json")
        
        existing_data = []
        if os.path.exists(filepath):
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    existing_data = json.load(f)
                    if not isinstance(existing_data, list):
                        existing_data = []
            except Exception as e:
                print(f"Warning: Could not read existing file {filepath}: {e}")
                existing_data = []
        
        existing_ids = {item.get('id') for item in existing_data if item.get('id')}
        
        new_items_to_add = []
        for item in new_data:
            item_id = item.get('id')
            if not item_id or item_id not in existing_ids:
                new_items_to_add.append(item)
                if item_id:
                    existing_ids.add(item_id)
        
        all_data = existing_data + new_items_to_add
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(all_data, f, indent=2, ensure_ascii=False)
        
        return filepath
    
    def _save_csv(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        if PANDAS_AVAILABLE:
            return self._save_csv_pandas(new_data, filename)
        else:
            return self._save_csv_manual(new_data, filename)
    
    def _save_csv_pandas(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        filepath = os.path.join(self.data_dir, f"{filename}.csv")
        
        new_df = pd.DataFrame(new_data)
        
        if os.path.exists(filepath):
            try:
                existing_df = pd.read_csv(filepath)
                combined_df = pd.concat([existing_df, new_df], ignore_index=True)
                
                if 'id' in combined_df.columns:
                    combined_df = combined_df.drop_duplicates(subset=['id'], keep='last')
                
            except Exception:
                combined_df = new_df
        else:
            combined_df = new_df
        
        combined_df.to_csv(filepath, index=False, encoding='utf-8')
        return filepath
    
    def _save_csv_manual(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        filepath = os.path.join(self.data_dir, f"{filename}.csv")
        
        all_fields = set()
        for item in new_data:
            all_fields.update(item.keys())
        
        existing_data = []
        if os.path.exists(filepath):
            try:
                with open(filepath, 'r', encoding='utf-8', newline='') as f:
                    reader = csv.DictReader(f)
                    existing_data = list(reader)
                    if reader.fieldnames:
                        all_fields.update(reader.fieldnames)
            except Exception as e:
                print(f"Warning: Could not read existing CSV file {filepath}: {e}")
                existing_data = []
        
        existing_ids = {item.get('id') for item in existing_data if item.get('id')}
        
        new_items_to_add = []
        for item in new_data:
            item_id = item.get('id')
            if not item_id or item_id not in existing_ids:
                new_items_to_add.append(item)
                if item_id:
                    existing_ids.add(item_id)
        
        all_data = existing_data + new_items_to_add
        
        with open(filepath, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=sorted(all_fields))
            writer.writeheader()
            writer.writerows(all_data)
        
        return filepath
    
    def _save_excel(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        if PANDAS_AVAILABLE:
            return self._save_excel_pandas(new_data, filename)
        else:
            csv_path = self._save_csv_manual(new_data, filename)
            return csv_path.replace('.csv', '_as_csv_instead_of_excel.csv')
    
    def _save_excel_pandas(self, new_data: List[Dict[str, Any]], filename: str) -> str:
        filepath = os.path.join(self.data_dir, f"{filename}.xlsx")
        
        new_df = pd.DataFrame(new_data)
        
        if os.path.exists(filepath):
            try:
                existing_df = pd.read_excel(filepath)
                combined_df = pd.concat([existing_df, new_df], ignore_index=True)
                
                if 'id' in combined_df.columns:
                    combined_df = combined_df.drop_duplicates(subset=['id'], keep='last')
                
            except Exception:
                combined_df = new_df
        else:
            combined_df = new_df
        
        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            combined_df.to_excel(writer, index=False, sheet_name='Scraped_Data')
        
        return filepath

class UniversalScraper:
    def __init__(self):
        if UA_AVAILABLE:
            self.ua = UserAgent()
        else:
            self.ua = SimpleFakeUserAgent()
        
        self.headers = {
            'User-Agent': self.ua.random,
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive',
        }
        
        if REQUESTS_AVAILABLE:
            self.session = requests.Session()
            self.session.headers.update(self.headers)
        else:
            self.session = None
    
    def scrape(self, url: str, max_items: int = 50, use_api: bool = False, platform: str = "auto") -> List[Dict[str, Any]]:
        try:
            if platform == "auto":
                platform = self._detect_platform(url)
            
            if 'reddit' in platform.lower():
                return self._scrape_reddit(url, max_items, use_api)
            elif 'linkedin' in platform.lower():
                return self._scrape_linkedin(url, max_items, use_api)
            else:
                return self._scrape_general(url, max_items)
        except Exception as e:
            st.error(f"Scraping failed: {e}")
            return []
    
    def _detect_platform(self, url: str) -> str:
        domain = urllib.parse.urlparse(url).netloc.lower()
        if 'reddit.com' in domain:
            return 'Reddit'
        elif 'linkedin.com' in domain:
            return 'LinkedIn'
        elif 'twitter.com' in domain or 'x.com' in domain:
            return 'Twitter/X'
        elif 'instagram.com' in domain:
            return 'Instagram'
        elif 'facebook.com' in domain:
            return 'Facebook'
        elif 'tiktok.com' in domain:
            return 'TikTok'
        elif 'youtube.com' in domain:
            return 'YouTube'
        else:
            return domain
    
    def _scrape_reddit(self, url: str, max_items: int, use_api: bool) -> List[Dict[str, Any]]:
        try:
            if not url.endswith('.json'):
                url = url.rstrip('/') + '.json'
            
            if REQUESTS_AVAILABLE:
                response = self.session.get(url, timeout=30)
                response.raise_for_status()
                data = response.json()
            else:
                req = urllib.request.Request(url, headers=self.headers)
                with urllib.request.urlopen(req, timeout=30) as response:
                    data = json.loads(response.read().decode('utf-8'))
            
            posts = []
            
            if isinstance(data, list) and len(data) > 0:
                post_data = data[0]['data']['children'][0]['data']
                posts.append({
                    'id': hashlib.md5(f"{post_data.get('id', '')}{url}".encode()).hexdigest()[:12],
                    'type': 'post',
                    'platform': 'Reddit',
                    'title': post_data.get('title', ''),
                    'content': post_data.get('selftext', ''),
                    'author': post_data.get('author', ''),
                    'url': post_data.get('url', ''),
                    'timestamp': post_data.get('created_utc', 0),
                    'engagement': post_data.get('score', 0),
                    'post_type': 'reddit_post'
                })
                
                if len(data) > 1 and 'data' in data[1]:
                    comments = data[1]['data']['children']
                    for comment in comments[:max_items-1]:
                        if comment.get('kind') == 't1':
                            comment_data = comment['data']
                            posts.append({
                                'id': hashlib.md5(f"{comment_data.get('id', '')}{url}".encode()).hexdigest()[:12],
                                'type': 'comment',
                                'platform': 'Reddit',
                                'title': '',
                                'content': comment_data.get('body', ''),
                                'author': comment_data.get('author', ''),
                                'url': '',
                                'timestamp': comment_data.get('created_utc', 0),
                                'engagement': comment_data.get('score', 0),
                                'post_type': 'reddit_comment'
                            })
            
            elif isinstance(data, dict) and 'data' in data:
                children = data['data'].get('children', [])
                for child in children[:max_items]:
                    if child.get('kind') == 't3':
                        post_data = child['data']
                        posts.append({
                            'id': hashlib.md5(f"{post_data.get('id', '')}{url}".encode()).hexdigest()[:12],
                            'type': 'post',
                            'platform': 'Reddit',
                            'title': post_data.get('title', ''),
                            'content': post_data.get('selftext', ''),
                            'author': post_data.get('author', ''),
                            'url': post_data.get('url', ''),
                            'timestamp': post_data.get('created_utc', 0),
                            'engagement': post_data.get('score', 0),
                            'post_type': 'reddit_post'
                        })
            
            return posts
            
        except Exception as e:
            st.error(f"Reddit scraping failed: {e}")
            return self._scrape_general(url, max_items)
    
    def _scrape_linkedin(self, url: str, max_items: int, use_api: bool) -> List[Dict[str, Any]]:
        return self._scrape_general(url, max_items)
    
    def _scrape_general(self, url: str, max_items: int) -> List[Dict[str, Any]]:
        try:
            if REQUESTS_AVAILABLE:
                response = self.session.get(url, timeout=30)
                response.raise_for_status()
                content = response.text
            else:
                req = urllib.request.Request(url, headers=self.headers)
                with urllib.request.urlopen(req, timeout=30) as response:
                    content = response.read().decode('utf-8', errors='ignore')
            
            if BS4_AVAILABLE:
                soup = BeautifulSoup(content, 'html.parser')
                
                for script in soup(["script", "style", "nav", "header", "footer"]):
                    script.decompose()
                
                posts = []
                post_selectors = [
                    'article', '[data-testid*="post"]', '[data-testid*="tweet"]',
                    '.post', '.entry', '.content', '.story', '.feed-item',
                    '[class*="post"]', '[class*="tweet"]', '[class*="story"]',
                    '[id*="post"]', '[id*="content"]'
                ]
                
                content_found = False
                for selector in post_selectors:
                    elements = soup.select(selector)
                    if elements:
                        for i, element in enumerate(elements[:max_items]):
                            text_content = element.get_text(separator=' ', strip=True)
                            if text_content and len(text_content) > 10:
                                posts.append({
                                    'id': hashlib.md5(f"{text_content}{url}{i}".encode()).hexdigest()[:12],
                                    'type': 'post',
                                    'platform': self._detect_platform(url),
                                    'title': self._extract_title(element),
                                    'content': text_content,
                                    'author': self._extract_author(element),
                                    'url': url,
                                    'timestamp': self._extract_timestamp(element),
                                    'engagement': self._extract_engagement(element),
                                    'post_type': f'scraped_content_{i+1}'
                                })
                        if posts:
                            content_found = True
                            break
                
                if not content_found:
                    paragraphs = soup.find_all(['p', 'div'], string=True)
                    for i, p in enumerate(paragraphs[:max_items]):
                        text = p.get_text().strip()
                        if len(text) > 20 and not self._is_navigation_text(text):
                            posts.append({
                                'id': hashlib.md5(f"{text}{url}{i}".encode()).hexdigest()[:12],
                                'type': 'text_block',
                                'platform': self._detect_platform(url),
                                'title': f'Text Block {i+1}',
                                'content': text,
                                'author': 'Unknown',
                                'url': url,
                                'timestamp': '',
                                'engagement': 0,
                                'post_type': 'paragraph_content'
                            })
                
                if posts:
                    return posts
            
            text = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            meaningful_lines = [
                line for line in lines 
                if len(line) > 30 and not re.match(r'^(home|about|contact|login|menu)', line.lower())
            ]
            
            posts = []
            domain = self._detect_platform(url)
            
            for i, line in enumerate(meaningful_lines[:max_items]):
                posts.append({
                    'id': hashlib.md5(f"{line}{url}{i}".encode()).hexdigest()[:12],
                    'type': 'extracted_text',
                    'platform': domain,
                    'title': f'Content {i+1}',
                    'content': line,
                    'author': 'Unknown',
                    'url': url,
                    'timestamp': '',
                    'engagement': 0,
                    'post_type': 'general_text'
                })
            
            return posts
            
        except Exception as e:
            st.error(f"General scraping failed: {e}")
            return []
    
    def _extract_title(self, element) -> str:
        if not BS4_AVAILABLE:
            return ''
        
        title_selectors = ['h1', 'h2', 'h3', '.title', '[class*="title"]', '[class*="header"]']
        for selector in title_selectors:
            title_elem = element.select_one(selector)
            if title_elem:
                return title_elem.get_text().strip()
        return ''
    
    def _extract_author(self, element) -> str:
        if not BS4_AVAILABLE:
            return 'Unknown'
        
        author_selectors = [
            '[class*="author"]', '[class*="user"]', '[class*="name"]',
            '[data-testid*="author"]', '[data-testid*="user"]'
        ]
        for selector in author_selectors:
            author_elem = element.select_one(selector)
            if author_elem:
                return author_elem.get_text().strip()
        return 'Unknown'
    
    def _extract_timestamp(self, element) -> str:
        if not BS4_AVAILABLE:
            return ''
        
        time_selectors = ['time', '[class*="time"]', '[class*="date"]']
        for selector in time_selectors:
            time_elem = element.select_one(selector)
            if time_elem:
                return time_elem.get('datetime') or time_elem.get_text().strip()
        return ''
    
    def _extract_engagement(self, element) -> int:
        if not BS4_AVAILABLE:
            return 0
        
        engagement_text = element.get_text()
        numbers = re.findall(r'\b\d+\b', engagement_text)
        return int(numbers[0]) if numbers else 0
    
    def _is_navigation_text(self, text: str) -> bool:
        nav_keywords = [
            'home', 'about', 'contact', 'privacy', 'terms', 'login', 'sign up',
            'menu', 'search', 'follow', 'share', 'like', 'comment', 'subscribe'
        ]
        return any(keyword in text.lower() for keyword in nav_keywords) and len(text) < 50

def main():
    st.set_page_config(
        page_title="NarrativeNexus: The Dynamic Text Analysis Platform",
        page_icon="🔍",
        layout="wide"
    )
    
    st.title(" NarrativeNexus: The Dynamic Text Analysis Platform")
    st.markdown("")
    
    # Main tabs for different input methods
    tab1, tab2 = st.tabs(["web scrap", "document pload"])
    
    with tab1:
        web_scraping_tab()
    
    with tab2:
        document_upload_tab()

def web_scraping_tab():
    with st.sidebar:
        st.header("Web Scraping Config")
        
        platform = st.selectbox(
            "Select Platform",
            ["Auto-Detect", "Reddit", "LinkedIn", "General/Other"],
            help="Choose the social media platform to scrape"
        )
        
        scraping_method = st.radio(
            "Scraping Method",
            ["Web Scraping", "API (if available)"],
            help="Choose between web scraping or API-based scraping"
        )
        
        use_api = scraping_method == "API (if available)"
        
        output_format = st.selectbox(
            "Output Format",
            ["JSON", "CSV", "Excel"],
            help="Choose the format for exported data"
        )
        
        custom_filename = st.text_input(
            "Custom Filename (optional)",
            value="social_media_data",
            help="Default: 'social_media_data' - new data will be appended to this file"
        )
        
        st.info("")
        
        if use_api:
            with st.expander("API Configuration"):
                st.info("API credentials can be configured here when available")
                st.text_input("API Key (optional)", type="password", help="Enter API key if required")
                st.text_input("Client ID (optional)", help="Enter client ID if required")
                st.text_input("Client Secret (optional)", type="password", help="Enter client secret if required")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("URL Input")
        url = st.text_input(
            "Enter URL to scrape:",
            placeholder="https://reddit.com/r/python or https://linkedin.com/posts/...",
            help="Enter the URL of the social media post or page you want to scrape"
        )
        
        with st.expander("Advanced Options"):
            max_items = st.number_input(
                "Maximum items to scrape",
                min_value=1,
                max_value=1000,
                value=50,
                help="Limit the number of items to scrape"
            )
            
            include_metadata = st.checkbox(
                "Include metadata",
                value=True,
                help="Include additional metadata like timestamps, URLs, etc."
            )
            
            st.subheader("Scraping Options")
            delay_between_requests = st.slider(
                "Delay between requests (seconds)",
                min_value=0.1,
                max_value=5.0,
                value=1.0,
                step=0.1,
                help="Add delay to be respectful to servers"
            )
            
            custom_headers = st.text_area(
                "Custom Headers (JSON format)",
                placeholder='{"User-Agent": "Custom Bot 1.0"}',
                help="Add custom headers for requests"
            )
    
    with col2:
        st.subheader("Quick Stats")
        if os.path.exists("data"):
            files = os.listdir("data")
            json_count = len([f for f in files if f.endswith('.json')])
            csv_count = len([f for f in files if f.endswith('.csv')])
            excel_count = len([f for f in files if f.endswith('.xlsx')])
            
            st.metric("JSON Files", json_count)
            st.metric("CSV Files", csv_count)
            st.metric("Excel Files", excel_count)
        else:
            st.info("No data files yet")
        
        st.subheader("Dependencies Status")
        st.write("✅ Streamlit")
        st.write("✅ Pandas" if PANDAS_AVAILABLE else "❌ Pandas")
        st.write("✅ Requests" if REQUESTS_AVAILABLE else "❌ Requests")
        st.write("✅ BeautifulSoup" if BS4_AVAILABLE else "❌ BeautifulSoup")
        st.write("✅ UserAgent" if UA_AVAILABLE else "❌ UserAgent")
        
        st.subheader("Document Support")
        st.write("✅ PDF" if (PYPDF2_AVAILABLE or PDFPLUMBER_AVAILABLE) else "❌ PDF")
        st.write("✅ DOCX" if DOCX_AVAILABLE else "❌ DOCX")
        st.write("✅ Excel" if (PANDAS_AVAILABLE or OPENPYXL_AVAILABLE) else "❌ Excel")
        st.write("✅ PowerPoint" if PPTX_AVAILABLE else "❌ PowerPoint")
        st.write("✅ Text/CSV")
    
    if st.button("🚀 Start Scraping", type="primary", use_container_width=True):
        if not url:
            st.error("Please enter a URL to scrape")
            return
        
        with st.spinner(f"Scraping data from {platform}..."):
            try:
                scraper = UniversalScraper()
                processor = DataProcessor()
                file_manager = FileManager()
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                status_text.text("Initializing scraper...")
                progress_bar.progress(20)
                
                platform_key = platform if platform != "Auto-Detect" else "auto"
                data = scraper.scrape(url, max_items=max_items, use_api=use_api, platform=platform_key)
                
                status_text.text("Processing data...")
                progress_bar.progress(60)
                
                processed_data = processor.process(data, include_metadata)
                
                status_text.text("Saving data...")
                progress_bar.progress(80)
                
                filename = file_manager.save_data(
                    processed_data,
                    output_format.lower(),
                    custom_filename or "social_media_data"
                )
                
                progress_bar.progress(100)
                status_text.text("Scraping completed!")
                
                display_results(processed_data, filename, output_format)
                
            except Exception as e:
                st.error(f"Error during scraping: {str(e)}")

def document_upload_tab():
    st.subheader("📄 Document Text Extraction")
    st.markdown("Upload documents to extract text content and add to your dataset")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        uploaded_files = st.file_uploader(
            "Choose files to upload",
            type=['pdf', 'docx', 'xlsx', 'xls', 'pptx', 'txt', 'csv'],
            accept_multiple_files=True,
            help="Supported formats: PDF, DOCX, XLSX, XLS, PPTX, TXT, CSV"
        )
        
        if uploaded_files:
            st.write(f"📁 **{len(uploaded_files)} file(s) selected:**")
            for file in uploaded_files:
                st.write(f"- {file.name} ({file.size:,} bytes)")
        
        with st.expander("Document Processing Options"):
            extract_by_page = st.checkbox(
                "Extract by page/sheet",
                value=True,
                help="Extract each page/sheet as separate entry (recommended)"
            )
            
            include_doc_metadata = st.checkbox(
                "Include document metadata",
                value=True,
                help="Include file information and extraction metadata"
            )
            
            merge_content = st.checkbox(
                "Merge all content",
                value=False,
                help="Combine all document content into single entries"
            )
    
    with col2:
        st.subheader("Document Support")
        doc_support = {
            "PDF": "✅" if (PYPDF2_AVAILABLE or PDFPLUMBER_AVAILABLE) else "⚠️",
            "DOCX": "✅" if DOCX_AVAILABLE else "⚠️", 
            "Excel": "✅" if (PANDAS_AVAILABLE or OPENPYXL_AVAILABLE) else "⚠️",
            "PowerPoint": "✅" if PPTX_AVAILABLE else "⚠️",
            "Text/CSV": "✅"
        }
        
        for format_name, status in doc_support.items():
            st.write(f"{status} {format_name}")
        
        if "⚠️" in doc_support.values():
            st.warning("Some formats may have limited support due to missing libraries")
        
        st.info("� Install additional packages for better support:\n- `pip install PyPDF2 pdfplumber`\n- `pip install python-docx`\n- `pip install python-pptx`\n- `pip install openpyxl xlrd`")
    
    col3, col4 = st.columns([1, 1])
    
    with col3:
        output_format_doc = st.selectbox(
            "Output Format",
            ["JSON", "CSV", "Excel"],
            key="doc_output_format",
            help="Choose the format for exported data"
        )
    
    with col4:
        custom_filename_doc = st.text_input(
            "Custom Filename (optional)",
            value="document_data",
            key="doc_filename",
            help="Default: 'document_data' - new data will be appended"
        )
    
    if st.button("📤 Extract Text from Documents", type="primary", use_container_width=True):
        if not uploaded_files:
            st.error("Please upload at least one file")
            return
        
        with st.spinner("Extracting text from documents..."):
            try:
                extractor = DocumentExtractor()
                processor = DataProcessor()
                file_manager = FileManager()
                
                all_extracted_data = []
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                for i, uploaded_file in enumerate(uploaded_files):
                    status_text.text(f"Processing {uploaded_file.name}...")
                    progress = (i + 1) / len(uploaded_files) * 0.8
                    progress_bar.progress(progress)
                    
                    file_content = uploaded_file.read()
                    extracted_data = extractor.extract_text(uploaded_file.name, file_content)
                    all_extracted_data.extend(extracted_data)
                
                status_text.text("Processing extracted data...")
                progress_bar.progress(0.9)
                
                processed_data = processor.process(all_extracted_data, include_doc_metadata)
                
                status_text.text("Saving data...")
                progress_bar.progress(0.95)
                
                filename = file_manager.save_data(
                    processed_data,
                    output_format_doc.lower(),
                    custom_filename_doc or "document_data"
                )
                
                progress_bar.progress(1.0)
                status_text.text("Document processing completed!")
                
                st.success(f"Successfully extracted text from {len(uploaded_files)} document(s)!")
                st.info(f"📄 Extracted {len(processed_data)} text blocks/pages")
                
                display_results(processed_data, filename, output_format_doc)
                
            except Exception as e:
                st.error(f"Error during document processing: {str(e)}")

def display_results(processed_data, filename, output_format):
    if os.path.exists(filename):
        st.info(f"✅ Data appended to existing file: {filename}")
        st.caption("💡 All previous data has been preserved")
    else:
        st.info(f"📁 New file created: {filename}")
        st.caption("💡 Future scrapes will be appended to this file")
    
    if processed_data:
        st.subheader("Data Preview")
        if PANDAS_AVAILABLE:
            df = pd.DataFrame(processed_data[:5])
            st.dataframe(df, use_container_width=True)
        else:
            for i, item in enumerate(processed_data[:5]):
                with st.expander(f"Item {i+1}: {item.get('title', 'Untitled')[:50]}"):
                    st.write(f"**Type:** {item.get('type', 'N/A')}")
                    st.write(f"**Platform:** {item.get('platform', 'N/A')}")
                    st.write(f"**Author:** {item.get('author', 'N/A')}")
                    st.write(f"**Content:** {item.get('content', 'N/A')[:200]}...")
        
        with open(filename, 'rb') as f:
            mime_types = {
                "json": "application/json",
                "csv": "text/csv",
                "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            }
            st.download_button(
                label=f"Download {output_format} file",
                data=f.read(),
                file_name=os.path.basename(filename),
                mime=mime_types.get(output_format.lower(), "application/octet-stream")
            )

if __name__ == "__main__":
    main()